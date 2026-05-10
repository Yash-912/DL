import os
from typing import List, Dict, Any


class DocumentParser:
    """
    Lightweight document parser for MVP.
    Uses direct parsers per format — no heavy ML models loaded at import time.
    Supported: PDF (.pdf), PowerPoint (.pptx), plain text (.txt, .md)
    """

    def parse(self, file_path: str) -> List[Dict[str, Any]]:
        """
        Parse a document and return a list of dicts with 'text' and 'metadata'.
        """
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"File not found: {file_path}")

        ext = os.path.splitext(file_path)[-1].lower()

        if ext == ".pdf":
            return self._parse_pdf(file_path)
        elif ext == ".pptx":
            return self._parse_pptx(file_path)
        elif ext in (".txt", ".md"):
            return self._parse_text(file_path)
        else:
            raise ValueError(f"Unsupported file type: {ext}")

    # ------------------------------------------------------------------
    # Private helpers
    # ------------------------------------------------------------------

    def _make_element(self, text: str, source_file: str, doc_type: str, page_number: int = 1) -> Dict[str, Any]:
        text = text.strip()
        if not text:
            return {}
        return {
            "text": text,
            "metadata": {
                "source_file": os.path.basename(source_file),
                "page_number": page_number,
                "doc_type": doc_type,
            },
        }

    def _parse_pdf(self, file_path: str) -> List[Dict[str, Any]]:
        import fitz  # PyMuPDF

        doc = fitz.open(file_path)
        results = []
        for page_num, page in enumerate(doc, start=1):
            text = page.get_text("text")
            elem = self._make_element(text, file_path, "pdf", page_number=page_num)
            if elem:
                results.append(elem)
        doc.close()
        return results

    def _parse_pptx(self, file_path: str) -> List[Dict[str, Any]]:
        from pptx import Presentation

        prs = Presentation(file_path)
        results = []
        for slide_num, slide in enumerate(prs.slides, start=1):
            parts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    parts.append(shape.text_frame.text)
            # Include speaker notes if present
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                parts.append(slide.notes_slide.notes_text_frame.text)
            full_text = "\n".join(parts)
            elem = self._make_element(full_text, file_path, "pptx", page_number=slide_num)
            if elem:
                results.append(elem)
        return results

    def _parse_text(self, file_path: str) -> List[Dict[str, Any]]:
        ext = os.path.splitext(file_path)[-1].lower().lstrip(".")
        with open(file_path, "r", encoding="utf-8") as f:
            content = f.read()
        elem = self._make_element(content, file_path, ext, page_number=1)
        return [elem] if elem else []
